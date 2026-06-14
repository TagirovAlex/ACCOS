import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def get_file_type(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
    if ext in image_exts:
        return "image"
    if ext == ".pdf":
        return "pdf"
    if ext == ".docx":
        return "docx"
    if ext == ".xlsx":
        return "xlsx"
    if ext == ".pptx":
        return "pptx"
    return "other"


def extract_text_from_pdf(file_path: str) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
    result = "\n".join(text_parts).strip()
    if not result:
        try:
            import fitz
            doc = fitz.open(file_path)
            text_parts = [page.get_text() for page in doc]
            result = "\n".join(text_parts).strip()
        except Exception:
            pass
    return result


def extract_text_from_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    rows = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows.append(f"=== {sheet} ===")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) if v is not None else "" for v in row]
            if any(v.strip() for v in vals):
                rows.append(" | ".join(vals))
    wb.close()
    return "\n".join(rows)


def extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    return "\n".join(text_parts)


def extract_text_from_image(file_path: str) -> str:
    from PIL import Image
    import pytesseract
    img = Image.open(file_path)
    return pytesseract.image_to_string(img, lang="rus+eng").strip()


def extract_text(file_path: str) -> str:
    file_type = get_file_type(file_path)
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    if file_type == "docx":
        return extract_text_from_docx(file_path)
    if file_type == "xlsx":
        return extract_text_from_xlsx(file_path)
    if file_type == "pptx":
        return extract_text_from_pptx(file_path)
    if file_type == "image":
        return extract_text_from_image(file_path)
    return ""


def get_file_size_limit(file_type: str) -> int:
    limits = {
        "image": 10 * 1024 * 1024,
        "pdf": 50 * 1024 * 1024,
        "docx": 10 * 1024 * 1024,
        "xlsx": 10 * 1024 * 1024,
        "pptx": 30 * 1024 * 1024,
    }
    return limits.get(file_type, 5 * 1024 * 1024)
