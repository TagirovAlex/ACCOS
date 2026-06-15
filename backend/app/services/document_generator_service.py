import logging
import re
from pathlib import Path

from app.core.paths import STATIC_DIR

logger = logging.getLogger(__name__)

DOCUMENTS_DIR = STATIC_DIR / "documents"


def _ensure_dirs():
    DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)


def generate_pdf(html_content: str, filename: str) -> str:
    from weasyprint import HTML
    _ensure_dirs()
    path = DOCUMENTS_DIR / f"{filename}.pdf"
    HTML(string=html_content).write_pdf(target=str(path))
    return str(path)


def generate_docx(template_path: str | None, fields: dict | None, content: str | None, filename: str) -> str:
    from docx import Document
    _ensure_dirs()
    path = DOCUMENTS_DIR / f"{filename}.docx"
    if template_path and Path(template_path).exists():
        doc = Document(template_path)
        if fields:
            for p in doc.paragraphs:
                for key, val in fields.items():
                    if f"{{{{{key}}}}}" in p.text:
                        p.text = p.text.replace(f"{{{{{key}}}}}", str(val))
        doc.save(str(path))
        return str(path)
    doc = Document()
    if content:
        for line in content.split("\n"):
            doc.add_paragraph(line)
    doc.save(str(path))
    return str(path)


def generate_xlsx(headers: list[str], rows: list[list], filename: str) -> str:
    import openpyxl
    from openpyxl.styles import Font
    _ensure_dirs()
    path = DOCUMENTS_DIR / f"{filename}.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    if headers:
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True)
    for row in rows:
        ws.append(row)
    wb.save(str(path))
    return str(path)


def generate_pptx(title: str, slides: list[dict], filename: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches
    _ensure_dirs()
    path = DOCUMENTS_DIR / f"{filename}.pptx"
    prs = Presentation()
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "")
        if slide_data.get("content"):
            slide.placeholders[1].text = slide_data["content"]
    prs.save(str(path))
    return str(path)


def apply_template(html: str, fields: dict) -> str:
    for key, val in fields.items():
        html = html.replace(f"{{{{{key}}}}}", str(val))
    return html


def generate(format: str, template_path: str | None = None, fields: dict | None = None,
             content: str | None = None, html: str | None = None,
             headers: list[str] | None = None, rows: list[list] | None = None,
             slides: list[dict] | None = None, title: str = "") -> str:
    import uuid
    filename = uuid.uuid4().hex
    if format == "pdf":
        if template_path and Path(template_path).exists():
            html_content = Path(template_path).read_text(encoding="utf-8")
            if fields:
                html_content = apply_template(html_content, fields)
            return generate_pdf(html_content, filename)
        return generate_pdf(html or content or "", filename)
    if format == "docx":
        return generate_docx(template_path, fields, content, filename)
    if format == "xlsx":
        return generate_xlsx(headers or [], rows or [], filename)
    if format == "pptx":
        return generate_pptx(title, slides or [], filename)
    raise ValueError(f"Unsupported format: {format}")
